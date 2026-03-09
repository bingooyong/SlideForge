from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Any, Callable, Dict, List, Optional, Tuple
import time
from concurrent.futures import ThreadPoolExecutor, as_completed

from pptx import Presentation

from app.config import settings
from app.core.logging import get_logger
from app.pipeline.background_cleaning import CleanedBackground, clean_background
from app.pipeline.color_extraction import StyleInfo, extract_colors
from app.pipeline.layout_ocr import (
    ImageBlock,
    Slide,
    TextBlock,
    analyze_layout,
    get_layout_ocr_raw_response_v2,
    parse_layout_v2_raw,
)
from app.pipeline.icon_cropping import CroppedImage, crop_imageblocks
from app.pipeline.pdf_to_images import PageImage, image_path_to_page_image, pdf_to_images
from app.pipeline.pptx_composition import add_slide_degraded, create_slide, setup_slide_size
from app.pipeline.pptx_composition_v2 import create_slide_v2
from app.pipeline.slide_schema_v2 import SlideDocumentV2, slide_document_v2_to_text_blocks_for_mask


@dataclass
class ProgressEvent:
    """
    Pipeline 进度事件，供外部（如 InMemoryTaskStore + WebSocket）消费。

    - task_id: 当前任务 ID
    - page_index: 当前处理的页索引（0-based）
    - total_pages: 总页数
    - stage: 当前处理阶段（pipeline 内部语义）
    - progress: 0–100 的整体进度百分比
    - message: 可选的人类可读说明
    """

    task_id: str
    page_index: int
    total_pages: int
    stage: str
    progress: int
    message: str = ""


ProgressCallback = Callable[[ProgressEvent], None]


@dataclass
class PipelinePageResult:
    """单页处理结果摘要。"""

    page_index: int
    success: bool
    error: Optional[str] = None


@dataclass
class PipelineResult:
    """
    整体流水线运行结果摘要。

    - output_pptx_path: 若至少一页成功且 PPTX 写入成功，则为输出文件路径；否则为 None
    - failed_pages: 处理失败的页索引列表（0-based）
    - page_results: 分页级别的详细结果
    """

    task_id: str
    output_pptx_path: Optional[Path]
    total_pages: int
    failed_pages: List[int]
    page_results: List[PipelinePageResult]


@dataclass
class _PageArtifacts:
    """单页中间产物，用于并发阶段与最终 PPTX 合成阶段解耦。"""

    page_index: int
    slide_schema: Slide
    style_info: StyleInfo
    text_blocks: List[TextBlock]
    image_blocks: List[ImageBlock]
    cropped_images: List[Optional[CroppedImage]]
    cleaned_bg: CleanedBackground
    slide_doc_v2: Optional[Any] = None  # V2 版面 doc，非空时用 create_slide_v2 合成


class PipelineService:
    """
    串联 3.1–3.6 模块的高阶服务。

    通过 `run()` 暴露单一入口，负责：
    - 调用 pdf_to_images 解析 PDF
    - 对每页依次执行：Gemini 版面解析 → 主题色提取 → 图标裁剪 → 背景净化 → PPTX 合成
    - 在各阶段通过 progress_callback 汇报进度
    - 对单页错误做降级处理并继续处理后续页面
    """

    # 单页内部阶段列表，用于计算整体进度
    _PAGE_STAGES: List[str] = [
        "ocr",
        "colors",
        "cropping",
        "background",
        "pptx",
    ]

    def __init__(self, default_output_root: Optional[Path] = None) -> None:
        """
        Args:
            default_output_root: 默认任务输出根目录，例如 Path("data/tasks")
        """
        self._default_output_root = default_output_root

    def run(
        self,
        task_id: str,
        pdf_path: str,
        *,
        mode: str = "standard",
        aspect_ratio: str = "16:9",
        output_root: Optional[Path] = None,
        progress_callback: Optional[ProgressCallback] = None,
        api_key: Optional[str] = None,
    ) -> PipelineResult:
        """
        执行完整流水线。

        Args:
            task_id: 任务 ID，用于构造输出路径与进度事件。
            pdf_path: 输入 PDF 本地路径。
            mode: 处理模式（当前占位，standard/lite 行为相同，后续 Task 可细化）。
            aspect_ratio: 输出 PPTX 页面宽高比（16:9 / 9:16 / 4:3）。
            output_root: 任务输出根目录，若为 None 则回退到初始化时的 default_output_root。
            progress_callback: 进度回调函数，可选。
            api_key: 可选，界面传入的 Gemini API Key；未传则使用环境 GEMINI_API_KEY。
        """
        logger = get_logger("pipeline.service")

        logger.info(
            "pipeline_task_start",
            task_id=task_id,
            pdf_path=pdf_path,
            mode=mode,
            aspect_ratio=aspect_ratio,
        )

        # 1) 解析输入为页面图片：PDF 多页渲染，图片单页直接加载
        if pdf_path.lower().endswith(".pdf"):
            pages = pdf_to_images(pdf_path)
        else:
            pages = [image_path_to_page_image(pdf_path)]
        total_pages = len(pages)

        # 任务目录：用于保存 output.pptx、每页输入图与 LLM 原始 JSON，便于与 run_pptx_from_raw 对比调试
        root = output_root or self._default_output_root
        task_dir: Optional[Path] = Path(root) / task_id if root else None
        if task_dir is not None:
            task_dir.mkdir(parents=True, exist_ok=True)
            for idx, page in enumerate(pages):
                (task_dir / f"page_{idx}_input.png").write_bytes(page.image_bytes)

        # 初始化 PPTX 文档
        prs = Presentation()
        setup_slide_size(prs, aspect_ratio)

        total_internal_stages = max(1, total_pages * len(self._PAGE_STAGES))

        def _notify(
            page_index: int,
            stage: str,
            stage_index: int,
            message: str = "",
        ) -> None:
            if not progress_callback:
                return
            completed_stage_count = page_index * len(self._PAGE_STAGES) + (stage_index + 1)
            progress = int(completed_stage_count / total_internal_stages * 100)
            event = ProgressEvent(
                task_id=task_id,
                page_index=page_index,
                total_pages=total_pages,
                stage=stage,
                progress=progress,
                message=message,
            )
            progress_callback(event)

        max_workers = max(1, min(getattr(settings, "PIPELINE_MAX_WORKERS", 3), total_pages))

        page_results_by_index: Dict[int, PipelinePageResult] = {}
        failed_pages: List[int] = []
        artifacts_by_page: Dict[int, _PageArtifacts] = {}

        def _process_page(page_idx: int, page: PageImage) -> Tuple[Optional[_PageArtifacts], PipelinePageResult]:
            page_start = time.perf_counter()
            logger.info(
                "pipeline_page_start",
                task_id=task_id,
                page=page_idx,
                total_pages=total_pages,
            )

            def _run_stage(
                stage_name: str,
                stage_index: int,
                func: Callable[[], object],
            ):
                stage_start = time.perf_counter()
                _notify(page_idx, stage_name, stage_index, f"Running {stage_name}")
                try:
                    result = func()
                    duration = time.perf_counter() - stage_start
                    logger.info(
                        "pipeline_stage_completed",
                        task_id=task_id,
                        page=page_idx,
                        stage=stage_name,
                        duration=duration,
                        result_status="success",
                    )
                    return result
                except Exception as exc:
                    duration = time.perf_counter() - stage_start
                    logger.error(
                        "pipeline_stage_error",
                        task_id=task_id,
                        page=page_idx,
                        stage=stage_name,
                        duration=duration,
                        error=str(exc),
                        result_status="failed",
                        exc_info=True,
                    )
                    raise

            try:
                # V2 流水线：先取 LLM 原始输出并落盘，再解析。解析结果与 scripts/run_pptx_from_raw.py 一致（同一份 raw 用脚本生成 PPTX 应与 Web 导出一致），便于以脚本为准调试。
                raw_response: str = _run_stage(
                    "ocr",
                    0,
                    lambda: get_layout_ocr_raw_response_v2(page, api_key=api_key),
                )
                if task_dir is not None:
                    raw_path = task_dir / f"page_{page_idx}_llm_raw_v2.txt"
                    raw_path.write_text(raw_response, encoding="utf-8")
                    logger.info(
                        "pipeline_raw_v2_saved",
                        task_id=task_id,
                        page=page_idx,
                        path=str(raw_path),
                    )

                slide_doc_v2: Any = parse_layout_v2_raw(raw_response)  # 与 run_pptx_from_raw 同源解析，不做 GLM 归一化

                if isinstance(slide_doc_v2, dict):
                    wrapper = SlideDocumentV2(
                        slide_metadata=None,
                        slide_data=slide_doc_v2,
                        elements=[],
                    )
                    text_blocks = slide_document_v2_to_text_blocks_for_mask(wrapper)
                else:
                    text_blocks = slide_document_v2_to_text_blocks_for_mask(slide_doc_v2)

                cleaned_bg: CleanedBackground = _run_stage(
                    "background",
                    3,
                    lambda: clean_background(
                        page_image=page,
                        text_blocks=text_blocks,
                    ),
                )

                # 占位：V2 路径不需要 colors/cropping 结果，但进度阶段 1、2 仍上报以保持进度条一致
                _run_stage(
                    "colors",
                    1,
                    lambda: StyleInfo(primaryColor="", backgroundColor="", accentColors=[]),
                )
                _run_stage("cropping", 2, lambda: crop_imageblocks(page_image=page, image_blocks=[]))

                page_duration = time.perf_counter() - page_start
                logger.info(
                    "pipeline_page_intermediate_completed",
                    task_id=task_id,
                    page=page_idx,
                    total_pages=total_pages,
                    duration=page_duration,
                    result_status="success",
                )

                # V2 用 slide_doc_v2 + cleaned_bg；slide_schema 等留空供降级或兼容
                artifacts = _PageArtifacts(
                    page_index=page_idx,
                    slide_schema=Slide(id="v2", index=page_idx, aspect_ratio="16:9", blocks=[]),
                    style_info=StyleInfo(primaryColor="", backgroundColor="", accentColors=[]),
                    text_blocks=text_blocks,
                    image_blocks=[],
                    cropped_images=[],
                    cleaned_bg=cleaned_bg,
                    slide_doc_v2=slide_doc_v2,
                )
                return artifacts, PipelinePageResult(page_index=page_idx, success=True, error=None)
            except Exception as exc:
                page_duration = time.perf_counter() - page_start
                logger.error(
                    "pipeline_page_failed",
                    task_id=task_id,
                    page=page_idx,
                    total_pages=total_pages,
                    duration=page_duration,
                    error=str(exc),
                    result_status="failed",
                    exc_info=True,
                )
                return None, PipelinePageResult(page_index=page_idx, success=False, error=str(exc))

        # 2) 页级并发：3.2–3.5 阶段并行执行
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            future_to_index = {
                executor.submit(_process_page, idx, page): idx
                for idx, page in enumerate(pages)
            }
            for future in as_completed(future_to_index):
                artifacts, result = future.result()
                page_results_by_index[result.page_index] = result
                if result.success and artifacts is not None:
                    artifacts_by_page[result.page_index] = artifacts
                else:
                    failed_pages.append(result.page_index)

        # 3) 串行 PPTX 合成（3.6）与最终 PPTX 写入
        for page_idx in range(total_pages):
            artifacts = artifacts_by_page.get(page_idx)
            if artifacts is None:
                _notify(page_idx, "pptx", 4, "Adding degraded slide")
                try:
                    add_slide_degraded(prs, pages[page_idx])
                except Exception as exc:
                    logger.error(
                        "pipeline_degraded_slide_error",
                        task_id=task_id,
                        page=page_idx,
                        error=str(exc),
                        exc_info=True,
                    )
                continue

            stage_start = time.perf_counter()
            _notify(page_idx, "pptx", 4, "Composing slide")
            try:
                if artifacts.slide_doc_v2 is not None:
                    create_slide_v2(
                        prs,
                        artifacts.slide_doc_v2,
                        artifacts.cleaned_bg,
                        icon_images=None,
                    )
                else:
                    create_slide(
                        prs=prs,
                        slide_schema=artifacts.slide_schema,
                        background=artifacts.cleaned_bg,
                        cropped_images=artifacts.cropped_images,
                        style_info=artifacts.style_info,
                    )
                duration = time.perf_counter() - stage_start
                logger.info(
                    "pipeline_stage_completed",
                    task_id=task_id,
                    page=page_idx,
                    stage="pptx",
                    duration=duration,
                    result_status="success",
                )
                # 覆盖先前的中间成功结果为最终成功状态
                page_results_by_index[page_idx] = PipelinePageResult(
                    page_index=page_idx,
                    success=True,
                    error=None,
                )
            except Exception as exc:
                duration = time.perf_counter() - stage_start
                logger.error(
                    "pipeline_stage_error",
                    task_id=task_id,
                    page=page_idx,
                    stage="pptx",
                    duration=duration,
                    error=str(exc),
                    result_status="failed",
                    exc_info=True,
                )
                failed_pages.append(page_idx)
                page_results_by_index[page_idx] = PipelinePageResult(
                    page_index=page_idx,
                    success=False,
                    error=str(exc),
                )

        page_results: List[PipelinePageResult] = [
            page_results_by_index[i]
            for i in sorted(page_results_by_index.keys())
        ]

        # 若至少一页成功，则写入 PPTX 文件（task_dir 已在 run 开头创建）
        output_pptx_path: Optional[Path] = None
        if any(r.success for r in page_results) and task_dir is not None:
            output_pptx_path = task_dir / "output.pptx"
            prs.save(str(output_pptx_path))

        result = PipelineResult(
            task_id=task_id,
            output_pptx_path=output_pptx_path,
            total_pages=total_pages,
            failed_pages=failed_pages,
            page_results=page_results,
        )
        logger.info(
            "pipeline_task_completed",
            task_id=task_id,
            total_pages=total_pages,
            failed_pages=failed_pages,
            has_output=output_pptx_path is not None,
        )
        return result

