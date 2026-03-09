"""
单页 PPTX 合成：将 Slide Schema、去字背景、裁剪图与 StyleInfo 合成为 python-pptx 单页。
供 Task 3.7 按页调用并合并为完整 PPTX。
"""

from __future__ import annotations

import io
from typing import List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from app.pipeline.background_cleaning import CleanedBackground
from app.pipeline.color_extraction import StyleInfo
from app.pipeline.layout_ocr_models import (
    ImageBlock,
    Slide,
    TextBlock,
)
from app.pipeline.icon_cropping import CroppedImage
from app.pipeline.pdf_to_images import PageImage

# 1 inch = 914400 EMU (python-pptx 内部)
EMU_PER_INCH = 914400.0


def _slide_dimensions_inch(prs: Presentation) -> tuple[float, float]:
    """返回 (width_inch, height_inch)。"""
    w_emu = prs.slide_width
    h_emu = prs.slide_height
    return w_emu / EMU_PER_INCH, h_emu / EMU_PER_INCH


def _apply_text_style(shape, text_block: TextBlock, style_info: StyleInfo) -> None:
    """对已创建的文本框 shape 应用 TextBlock 与 StyleInfo 的样式。"""
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
    p.text = text_block.content or ""

    font = p.font
    if text_block.style:
        font.size = Pt(max(1, int(text_block.style.fontSize)))
        if text_block.style.fontWeight and "bold" in text_block.style.fontWeight.lower():
            font.bold = True
        if text_block.style.fontColor:
            try:
                hex_color = text_block.style.fontColor.lstrip("#")
                r = int(hex_color[0:2], 16)
                g = int(hex_color[2:4], 16)
                b = int(hex_color[4:6], 16)
                font.color.rgb = RGBColor(r, g, b)
            except (ValueError, IndexError):
                font.color.rgb = RGBColor(0, 0, 0)
        if text_block.style.textAlign:
            from pptx.enum.text import PP_ALIGN

            align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}
            p.alignment = align_map.get(text_block.style.textAlign, PP_ALIGN.LEFT)
    else:
        font.size = Pt(12)
        font.color.rgb = RGBColor(0, 0, 0)


def create_slide(
    prs: Presentation,
    slide_schema: Slide,
    background: CleanedBackground,
    cropped_images: List[Optional[CroppedImage]],
    style_info: StyleInfo,
) -> object:
    """
    在给定的 Presentation 中追加一页幻灯片：背景层使用去字图，按 Schema 插入图片块与文本框。

    假定 prs 的 slide_width / slide_height 已由调用方按目标比例设置（如 16:9/4:3/9:16）。
    cropped_images 与 slide_schema.blocks 中的 ImageBlock 顺序一致，缺失或失败处为 None。

    Returns:
        python-pptx 的 Slide 对象（已加入 prs.slides）。
    """
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)

    w_inch, h_inch = _slide_dimensions_inch(prs)

    # 1) 背景层：全页去字图
    slide.shapes.add_picture(
        io.BytesIO(background.image_bytes),
        Inches(0),
        Inches(0),
        width=Inches(w_inch),
        height=Inches(h_inch),
    )

    # 2) 图片块：按 ImageBlock 顺序与 cropped_images 一一对应
    image_index = 0
    for block in slide_schema.blocks:
        if not isinstance(block, ImageBlock):
            continue
        if image_index >= len(cropped_images):
            break
        cropped = cropped_images[image_index]
        image_index += 1
        if cropped is None:
            continue

        box = block.box
        left = Inches(box.x * w_inch)
        top = Inches(box.y * h_inch)
        width = Inches(box.w * w_inch)
        height = Inches(box.h * h_inch)
        if width <= 0 or height <= 0:
            continue

        try:
            slide.shapes.add_picture(
                io.BytesIO(cropped.image_bytes),
                left,
                top,
                width=width,
                height=height,
            )
        except Exception:
            continue

    # 3) 文本框：按 bbox 与样式
    for block in slide_schema.blocks:
        if not isinstance(block, TextBlock):
            continue
        box = block.box
        left = Inches(box.x * w_inch)
        top = Inches(box.y * h_inch)
        width = Inches(box.w * w_inch)
        height = Inches(box.h * h_inch)
        if width <= 0 or height <= 0:
            continue

        shape = slide.shapes.add_textbox(left, top, width, height)
        _apply_text_style(shape, block, style_info)

    return slide


def add_slide_degraded(prs: Presentation, page_image: PageImage) -> object:
    """
    追加一页降级占位幻灯片：仅使用原页图整页铺满，无文字、无去字。
    用于单页 pipeline 失败时保持 PPT 页数与 PDF 一致。
    """
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    w_inch, h_inch = _slide_dimensions_inch(prs)
    slide.shapes.add_picture(
        io.BytesIO(page_image.image_bytes),
        Inches(0),
        Inches(0),
        width=Inches(w_inch),
        height=Inches(h_inch),
    )
    return slide


def slide_size_for_aspect_ratio(aspect_ratio: str) -> tuple[float, float]:
    """
    返回 (width_inch, height_inch)，用于在创建 Presentation 时设置 slide 尺寸。
    支持 16:9, 4:3, 9:16。
    """
    if aspect_ratio == "16:9":
        return (13.333, 7.5)
    if aspect_ratio == "4:3":
        return (10.0, 7.5)
    if aspect_ratio == "9:16":
        return (7.5, 13.333)
    return (13.333, 7.5)


def setup_slide_size(prs: Presentation, aspect_ratio: str) -> None:
    """在添加幻灯片前设置演示文稿的页面尺寸（与前端 aspectRatio 一致）。"""
    w_inch, h_inch = slide_size_for_aspect_ratio(aspect_ratio)
    prs.slide_width = int(w_inch * EMU_PER_INCH)
    prs.slide_height = int(h_inch * EMU_PER_INCH)
