"""
单页 PPTX 合成 V2：按「容器嵌套 + 相对坐标」Schema V2 渲染。

支持 group、shape_text_box、icon_text_layout、list_block 等，
实现圆角容器、富文本分段、图标+文字混排、列表。
"""

from __future__ import annotations

import io
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

from app.pipeline.background_cleaning import CleanedBackground
from app.pipeline.default_icon_images import build_default_icon_images, resolve_icon
from app.pipeline.slide_schema_v2 import Bbox

EMU_PER_INCH = 914400.0
DEFAULT_FONT_FAMILY = "Microsoft YaHei"
DEFAULT_TEXT_PADDING_PT = 4
TITLE_TEXT_PADDING_PT = 6
DEFAULT_LINE_SPACING = 1.18
DEFAULT_SPACE_AFTER_PT = 2

# 圆角半径：python-pptx 用 adjustments[0] 控制，约 0.04 为轻微圆角，默认 0.16667 偏大
DEFAULT_CORNER_ADJUSTMENT = 0.04
MAX_CORNER_ADJUSTMENT = 0.08


def _slide_dimensions_inch(prs: Presentation) -> tuple[float, float]:
    w_emu = prs.slide_width
    h_emu = prs.slide_height
    return w_emu / EMU_PER_INCH, h_emu / EMU_PER_INCH


def _hex_to_rgb(hex_color: Optional[str]) -> Optional[RGBColor]:
    if not hex_color:
        return None
    hex_color = hex_color.lstrip("#")
    if len(hex_color) != 6:
        return None
    try:
        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)
        return RGBColor(r, g, b)
    except ValueError:
        return None


def _radius_to_adjustment(radius: Optional[Any]) -> float:
    """
    将 JSON 的 radius 映射为 python-pptx adjustments[0]。
    - 若为小数（如 0.008、0.015），视为相对比例，线性映射到 DEFAULT~MAX；
    - 若为整数（如 6、12），按原逻辑除以 200。
    """
    if radius is None:
        return DEFAULT_CORNER_ADJUSTMENT
    try:
        r = float(radius)
    except (TypeError, ValueError):
        return DEFAULT_CORNER_ADJUSTMENT
    if 0 < r < 1:
        return min(MAX_CORNER_ADJUSTMENT, max(DEFAULT_CORNER_ADJUSTMENT, r * 2.0))
    if r >= 1:
        return min(MAX_CORNER_ADJUSTMENT, r / 200.0)
    return DEFAULT_CORNER_ADJUSTMENT


def relative_bbox_to_absolute(
    relative_bbox: Bbox,
    parent_absolute_bbox: tuple[float, float, float, float],
) -> tuple[float, float, float, float]:
    """
    将子元素的相对坐标 [x, y, w, h]（0~1）转为幻灯片/画布上的绝对坐标（英寸）。
    父容器绝对 bbox 为 (p_x, p_y, p_w, p_h)。
    渲染时 _render_element 内已按此逻辑递归换算。
    """
    if len(relative_bbox) != 4:
        return 0.0, 0.0, 0.0, 0.0
    p_x, p_y, p_w, p_h = parent_absolute_bbox
    r_x, r_y, r_w, r_h = relative_bbox
    abs_x = p_x + r_x * p_w
    abs_y = p_y + r_y * p_h
    abs_w = r_w * p_w
    abs_h = r_h * p_h
    return (abs_x, abs_y, abs_w, abs_h)


def _get_attr(obj: Any, key: str, default: Any = None) -> Any:
    """从 Pydantic 或 dict 取属性。"""
    if isinstance(obj, dict):
        return obj.get(key, default)
    return getattr(obj, key, default)


def _get_runs_from_elem(elem: Any) -> List[Any]:
    """
    从元素中统一取出 text runs 列表，兼容多种大模型输出格式（页面导出与命令行共用）：
    text_runs、content（数组）、content.runs（对象）、text_content、text_box.text_runs。
    """
    runs = _get_attr(elem, "text_runs") or _get_attr(elem, "text_content") or _get_attr(elem, "text")
    if isinstance(runs, list):
        return runs
    content = _get_attr(elem, "content")
    if isinstance(content, list):
        return content
    if isinstance(content, dict) and isinstance(content.get("runs"), list):
        return content["runs"]
    text_box = _get_attr(elem, "text_box")
    if isinstance(text_box, dict):
        runs = text_box.get("text_runs") or (isinstance(text_box.get("content"), dict) and text_box.get("content", {}).get("runs"))
        if isinstance(runs, list):
            return runs
    return []


def _run_font_size_pt(run_spec: Any, slide_h_inch: float) -> Optional[float]:
    """
    从 run 中解析字号并转为 Pt。
    若为相对比例（0~1，如 0.028），则按幻灯片高度换算：h_inch * 72 * ratio。
    兼容 font_size 与 font_size_ratio 字段。
    """
    fs = None
    if isinstance(run_spec, dict):
        fs = run_spec.get("font_size") or run_spec.get("font_size_ratio")
    elif hasattr(run_spec, "style") and run_spec.style:
        fs = getattr(run_spec.style, "font_size", None)
    if fs is None:
        return None
    try:
        v = float(fs)
    except (TypeError, ValueError):
        return None
    if 0 < v < 1:
        return max(8.0, min(72.0, slide_h_inch * 72.0 * v))
    return max(8.0, min(72.0, v))


def _run_color(run_spec: Any) -> Optional[str]:
    """从 run 取颜色（color / font_color 或 style.color）。"""
    if isinstance(run_spec, dict):
        return run_spec.get("color") or run_spec.get("font_color")
    if hasattr(run_spec, "style") and run_spec.style:
        return getattr(run_spec.style, "color", None)
    return getattr(run_spec, "color", None) or getattr(run_spec, "font_color", None)


def _run_bold(run_spec: Any) -> bool:
    """从 run 取是否粗体。"""
    fw = None
    if isinstance(run_spec, dict):
        fw = run_spec.get("font_weight")
    elif hasattr(run_spec, "style") and run_spec.style:
        fw = getattr(run_spec.style, "font_weight", None)
    return bool(fw and "bold" in (fw or "").lower())


def _run_text(run_spec: Any) -> str:
    """从 run 取文本。"""
    if isinstance(run_spec, dict):
        return run_spec.get("text") or ""
    return getattr(run_spec, "text", None) or ""


def _set_cjk_font(run: Any, font_name: str = DEFAULT_FONT_FAMILY) -> None:
    """为 run 同时设置西文和东亚字体，提升中文粗体兼容性。"""
    try:
        run.font.name = font_name
    except Exception:
        pass
    try:
        r_pr = run._r.get_or_add_rPr()
        r_pr.set(qn("a:ea"), font_name)
    except Exception:
        pass


def _apply_run_style(run: Any, run_spec: Any, slide_h_inch: float) -> None:
    pt_size = _run_font_size_pt(run_spec, slide_h_inch)
    if pt_size is not None:
        run.font.size = Pt(pt_size)
    _set_cjk_font(run)
    if _run_bold(run_spec):
        run.font.bold = True
    rgb = _hex_to_rgb(_run_color(run_spec))
    if rgb:
        run.font.color.rgb = rgb


def _inflate_text_region(
    left_inch: float,
    top_inch: float,
    width_inch: float,
    height_inch: float,
    parent_abs: tuple[float, float, float, float],
    x_expand_ratio: float = 0.04,
    y_expand_ratio: float = 0.06,
) -> tuple[float, float, float, float]:
    """对文本框做轻微安全扩展，避免 bbox 过窄导致异常换行。"""
    if width_inch <= 0 or height_inch <= 0:
        return left_inch, top_inch, width_inch, height_inch
    px, py, pw, ph = parent_abs

    expand_x = width_inch * x_expand_ratio / 2.0
    expand_y = height_inch * y_expand_ratio / 2.0

    new_left = max(px, left_inch - expand_x)
    new_top = max(py, top_inch - expand_y)
    max_right = px + pw
    max_bottom = py + ph
    new_right = min(max_right, left_inch + width_inch + expand_x)
    new_bottom = min(max_bottom, top_inch + height_inch + expand_y)

    return new_left, new_top, max(0.0, new_right - new_left), max(0.0, new_bottom - new_top)


_ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}


def _parse_align_string(align: Optional[str]) -> Optional[PP_ALIGN]:
    """将元素级 text_align / alignment 字符串转为 PP_ALIGN。"""
    if not align:
        return None
    return _ALIGN_MAP.get(str(align).strip().lower(), None)


def _align_from_run(run_spec: Any) -> Optional[PP_ALIGN]:
    align = None
    if isinstance(run_spec, dict):
        align = run_spec.get("align")
    elif hasattr(run_spec, "style") and run_spec.style:
        align = getattr(run_spec.style, "align", None)
    if not align:
        return None
    return _ALIGN_MAP.get(str(align).lower(), PP_ALIGN.LEFT)


def _apply_text_runs(
    shape: Any,
    runs: List[Any],
    slide_h_inch: float = 5.625,
    padding_pt: int = DEFAULT_TEXT_PADDING_PT,
    default_align: Optional[str] = None,
    vertical_anchor: Optional[Any] = None,
) -> None:
    """
    向 shape.text_frame 按 runs 写入富文本。
    支持元素级 text_align/alignment（default_align），run 内 align 可覆盖。
    支持 vertical_anchor 实现垂直居中（如 shape_text_box 标题）。
    """
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    if vertical_anchor is not None:
        tf.vertical_anchor = vertical_anchor

    try:
        tf.margin_left = Pt(padding_pt)
        tf.margin_right = Pt(padding_pt)
        tf.margin_top = Pt(1)
        tf.margin_bottom = Pt(1)
    except Exception:
        pass

    p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
    p.clear()
    p.line_spacing = DEFAULT_LINE_SPACING
    p.space_after = Pt(DEFAULT_SPACE_AFTER_PT)

    elem_align = _parse_align_string(default_align)
    if elem_align is not None:
        p.alignment = elem_align

    if not runs:
        p.add_run().text = " "
        return

    for run_spec in runs:
        run = p.add_run()
        run.text = _run_text(run_spec)
        _apply_run_style(run, run_spec, slide_h_inch)
        run_align = _align_from_run(run_spec)
        if run_align is not None:
            p.alignment = run_align


def _force_bullet_paragraph(paragraph: Any, bullet_char: str = "•") -> None:
    """强制段落开启项目符号，避免模板无默认 bullet 样式时丢失圆点。"""
    try:
        p_pr = paragraph._p.get_or_add_pPr()
        for tag in ("a:buNone", "a:buAutoNum", "a:buBlip", "a:buChar"):
            node = p_pr.find(qn(tag))
            if node is not None:
                p_pr.remove(node)
        bu_char = OxmlElement("a:buChar")
        bu_char.set("char", bullet_char)
        p_pr.append(bu_char)
    except Exception:
        pass


def _add_soft_shadow(
    slide: Any,
    left_inch: float,
    top_inch: float,
    width_inch: float,
    height_inch: float,
    radius: Optional[Any],
    offset_pt: float = 2.0,
    transparency: float = 0.84,
) -> None:
    """通过底层偏移圆角形状模拟柔和阴影（python-pptx 阴影 API 能力有限）。"""
    dx = offset_pt / 72.0
    dy = offset_pt / 72.0
    sp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left_inch + dx),
        Inches(top_inch + dy),
        Inches(width_inch),
        Inches(height_inch),
    )
    if getattr(sp, "adjustments", None) and len(sp.adjustments) > 0:
        try:
            sp.adjustments[0] = _radius_to_adjustment(radius)
        except Exception:
            pass
    sp.fill.solid()
    sp.fill.fore_color.rgb = RGBColor(45, 59, 73)
    try:
        sp.fill.transparency = transparency
    except Exception:
        pass
    sp.line.fill.background()


def _add_rounded_rect(
    slide: Any,
    left_inch: float,
    top_inch: float,
    width_inch: float,
    height_inch: float,
    fill_color: Optional[str] = None,
    border_color: Optional[str] = None,
    border_width: Optional[int] = None,
    shadow: Optional[bool] = None,
    radius: Optional[int] = None,
    shape_type: Optional[str] = None,
) -> Any:
    use_rounded = shape_type != "rectangle"

    # 白色大卡片默认加柔和阴影，或当 JSON 显式要求 shadow 时加阴影；有边框的容器（如红框）绝不自动加阴影。
    should_add_shadow = shadow is True
    if not should_add_shadow and fill_color and fill_color.lower() in {"#fff", "#ffffff"}:
        if not border_color:
            should_add_shadow = width_inch >= 2.2 and height_inch >= 1.2

    if should_add_shadow:
        _add_soft_shadow(
            slide,
            left_inch,
            top_inch,
            width_inch,
            height_inch,
            radius=radius,
        )

    sp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE if not use_rounded else MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left_inch),
        Inches(top_inch),
        Inches(width_inch),
        Inches(height_inch),
    )
    if use_rounded and getattr(sp, "adjustments", None) and len(sp.adjustments) > 0:
        try:
            sp.adjustments[0] = _radius_to_adjustment(radius)
        except Exception:
            sp.adjustments[0] = DEFAULT_CORNER_ADJUSTMENT

    # 仅当 JSON 明确指定 shadow: true 时保留 ppt 原生阴影，否则关闭。
    if shadow is not True and getattr(sp, "shadow", None) is not None:
        try:
            sp.shadow.inherit = False
        except Exception:
            pass

    sp.fill.solid()
    if fill_color:
        rgb = _hex_to_rgb(fill_color)
        if rgb:
            sp.fill.fore_color.rgb = rgb
    else:
        sp.fill.fore_color.rgb = RGBColor(255, 255, 255)

    if border_color and border_width:
        sp.line.color.rgb = _hex_to_rgb(border_color) or RGBColor(0, 0, 0)
        sp.line.width = Pt(max(1, border_width))
    else:
        sp.line.fill.background()
    return sp


def _extract_background_spec(doc: Any) -> Optional[dict]:
    """提取 slide_metadata.background（兼容 dict 与 Pydantic）。"""
    if isinstance(doc, dict):
        slide_metadata = doc.get("slide_metadata")
        if isinstance(slide_metadata, dict):
            bg = slide_metadata.get("background")
            return bg if isinstance(bg, dict) else None
        return None

    metadata = getattr(doc, "slide_metadata", None)
    if metadata is None:
        return None
    bg = getattr(metadata, "background", None)
    if bg is None:
        return None
    if isinstance(bg, dict):
        return bg

    return {
        "type": getattr(bg, "type", None),
        "start_color": getattr(bg, "start_color", None),
        "end_color": getattr(bg, "end_color", None),
        "direction": getattr(bg, "direction", None),
        "colors": getattr(bg, "colors", None),
    }


def _try_draw_gradient_background(
    slide: Any,
    w_inch: float,
    h_inch: float,
    bg_spec: Optional[dict],
) -> bool:
    """尝试绘制全屏渐变背景，成功返回 True。"""
    if not bg_spec or str(bg_spec.get("type") or "").lower() != "gradient":
        return False

    start_color = bg_spec.get("start_color")
    end_color = bg_spec.get("end_color")
    colors = bg_spec.get("colors") or []
    if not start_color and isinstance(colors, list) and colors:
        start_color = colors[0]
    if not end_color and isinstance(colors, list) and len(colors) >= 2:
        end_color = colors[-1]

    start_rgb = _hex_to_rgb(start_color)
    end_rgb = _hex_to_rgb(end_color)
    if not start_rgb:
        return False
    if not end_rgb:
        end_rgb = start_rgb

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(w_inch),
        Inches(h_inch),
    )
    shape.line.fill.background()

    fill = shape.fill
    try:
        fill.gradient()
        gradient_stops = fill.gradient_stops
        if len(gradient_stops) >= 2:
            gradient_stops[0].position = 0.0
            gradient_stops[0].color.rgb = start_rgb
            gradient_stops[1].position = 1.0
            gradient_stops[1].color.rgb = end_rgb
            direction = str(bg_spec.get("direction") or "").lower()
            if direction in {"vertical", "top_to_bottom"}:
                fill.gradient_angle = 90
            elif direction in {"horizontal", "left_to_right"}:
                fill.gradient_angle = 0
        else:
            fill.solid()
            fill.fore_color.rgb = start_rgb
    except Exception:
        fill.solid()
        fill.fore_color.rgb = start_rgb
    return True


def _render_element(
    slide: Any,
    elem: Any,
    w_inch: float,
    h_inch: float,
    parent_origin: tuple[float, float],
    parent_size: tuple[float, float],
    icon_images: Dict[str, bytes],
) -> None:
    """
    在 slide 上渲染单个元素；支持 Pydantic 或 dict。
    子元素 bbox 为相对父容器的 0~1 比例，递归换算为绝对坐标。
    支持 type: text_box / text_block / group / shape_text_box / icon_text_layout / list_block / list_layout。
    """
    px, py = parent_origin
    pw, ph = parent_size
    b = _get_attr(elem, "bbox")
    if not b or len(b) != 4:
        return
    left_inch, top_inch, width_inch, height_inch = relative_bbox_to_absolute(b, (px, py, pw, ph))
    if width_inch <= 0 or height_inch <= 0:
        return

    elem_type = _get_attr(elem, "type") or ""
    parent_abs = (px, py, pw, ph)

    # ----- 带 list_items 的 text_box（JSON 中红框内三条列表等）按列表渲染，否则会漏字 -----
    if elem_type in ("text_box", "text_block"):
        list_items = _get_attr(elem, "list_items") or []
        if list_items:
            items = list_items
            tx, ty, tw, th = _inflate_text_region(left_inch, top_inch, width_inch, height_inch, parent_abs)
            shape = slide.shapes.add_textbox(Inches(tx), Inches(ty), Inches(tw), Inches(th))
            tf = shape.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.NONE
            try:
                tf.margin_left = Pt(DEFAULT_TEXT_PADDING_PT)
                tf.margin_right = Pt(DEFAULT_TEXT_PADDING_PT)
                tf.margin_top = Pt(2)
                tf.margin_bottom = Pt(2)
            except Exception:
                pass
            bullet_type = (_get_attr(elem, "bullet_type") or "disc").lower()
            bullet_char = "•" if bullet_type in {"disc", "dot", "bullet"} else "-"
            for i, item in enumerate(items):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                if hasattr(p, "clear"):
                    p.clear()
                p.level = 0
                p.line_spacing = DEFAULT_LINE_SPACING
                p.space_after = Pt(DEFAULT_SPACE_AFTER_PT)
                _force_bullet_paragraph(p, bullet_char=bullet_char)
                if isinstance(item, list):
                    runs = item
                else:
                    runs = _get_attr(item, "text_runs") or _get_attr(item, "runs")
                    if not isinstance(runs, list) and isinstance(_get_attr(item, "content"), dict):
                        runs = _get_attr(_get_attr(item, "content"), "runs") or []
                    if not isinstance(runs, list):
                        runs = _get_attr(item, "text") or _get_attr(item, "content") or []
                    runs = runs if isinstance(runs, list) else []
                for run_spec in runs:
                    run = p.add_run()
                    run.text = _run_text(run_spec)
                    _apply_run_style(run, run_spec, h_inch)
            return

    # ----- text_box 或 text_block（无 list_items）-----
    if elem_type in ("text_box", "text_block"):
        runs = _get_runs_from_elem(elem)
        tx, ty, tw, th = _inflate_text_region(left_inch, top_inch, width_inch, height_inch, parent_abs)
        shape = slide.shapes.add_textbox(Inches(tx), Inches(ty), Inches(tw), Inches(th))
        elem_align = _get_attr(elem, "text_align") or _get_attr(elem, "alignment")
        _apply_text_runs(shape, runs, h_inch, default_align=elem_align)
        return

    # ----- shape_text_box -----
    if elem_type == "shape_text_box":
        fill = (
            _get_attr(elem, "fill_color")
            or _get_attr(elem, "background_color")
            or (_get_attr(elem, "shape") and _get_attr(_get_attr(elem, "shape"), "fill_color"))
        )
        if fill:
            _add_rounded_rect(
                slide,
                left_inch,
                top_inch,
                width_inch,
                height_inch,
                fill_color=fill,
                border_color=_get_attr(elem, "border_color")
                or (
                    _get_attr(elem, "shape")
                    and _get_attr(_get_attr(elem, "shape"), "border")
                    and _get_attr(_get_attr(_get_attr(elem, "shape"), "border"), "color")
                ),
                border_width=_get_attr(elem, "border_width")
                or (
                    _get_attr(elem, "shape")
                    and _get_attr(_get_attr(elem, "shape"), "border")
                    and _get_attr(_get_attr(_get_attr(elem, "shape"), "border"), "width")
                ),
                radius=_get_attr(elem, "border_radius")
                or _get_attr(elem, "radius")
                or (_get_attr(elem, "shape") and _get_attr(_get_attr(elem, "shape"), "radius")),
                shape_type=_get_attr(elem, "shape") and _get_attr(_get_attr(elem, "shape"), "type"),
            )

        runs = _get_runs_from_elem(elem)
        tx, ty, tw, th = _inflate_text_region(left_inch, top_inch, width_inch, height_inch, parent_abs)
        shape = slide.shapes.add_textbox(Inches(tx), Inches(ty), Inches(tw), Inches(th))
        elem_align = _get_attr(elem, "text_align") or _get_attr(elem, "alignment")
        _apply_text_runs(
            shape, runs, h_inch,
            padding_pt=TITLE_TEXT_PADDING_PT,
            default_align=elem_align,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
        )
        return

    # ----- icon_text_layout：兼容 icon/icon_properties、text_bbox/text_content_bbox/text_box.bbox、runs 多种格式 -----
    if elem_type == "icon_text_layout":
        icon = _get_attr(elem, "icon") or _get_attr(elem, "icon_properties")
        pid = _get_attr(elem, "placeholder_id") or (icon and _get_attr(icon, "placeholder_id"))
        icon_bbox = _get_attr(elem, "icon_bbox") or (icon and _get_attr(icon, "bbox"))
        text_box = _get_attr(elem, "text_box")
        text_bbox = (
            _get_attr(elem, "text_bbox")
            or _get_attr(elem, "text_content_bbox")
            or (text_box and _get_attr(text_box, "bbox"))
        )
        runs = _get_runs_from_elem(elem)
        current_abs = (left_inch, top_inch, width_inch, height_inch)

        actual_iw = 0
        ix, iy = 0.0, 0.0
        if icon_bbox and len(icon_bbox) == 4 and pid:
            icon_bytes = resolve_icon(pid, icon_images)
            if icon_bytes:
                ix, iy, iw, ih = relative_bbox_to_absolute(icon_bbox, current_abs)

                # FIX 1: 限制图标最大宽度（0.35 英寸），防止大模型估算过大变成巨无霸
                max_icon_inch = 0.35
                actual_iw = min(iw, max_icon_inch)

                if actual_iw > 0:
                    try:
                        slide.shapes.add_picture(
                            io.BytesIO(icon_bytes),
                            Inches(ix),
                            Inches(iy),
                            width=Inches(actual_iw),
                        )
                    except Exception:
                        pass

        if text_bbox and len(text_bbox) == 4:
            tx, ty, tw, th = relative_bbox_to_absolute(text_bbox, current_abs)
        else:
            tx, ty, tw, th = left_inch, top_inch, width_inch, height_inch

        # FIX 2: 强行接管图文间距！如果存在图标，强制文字紧贴在图标右侧 0.1 英寸处
        if actual_iw > 0:
            expected_tx = ix + actual_iw + 0.1
            if tx > expected_tx:
                diff = tx - expected_tx
                tx = expected_tx
                tw = tw + diff

        tx, ty, tw, th = _inflate_text_region(tx, ty, tw, th, current_abs)

        if tw > 0 and th > 0:
            shape = slide.shapes.add_textbox(Inches(tx), Inches(ty), Inches(tw), Inches(th))
            elem_align = _get_attr(elem, "text_align") or _get_attr(elem, "alignment")
            _apply_text_runs(shape, runs, h_inch, default_align=elem_align)
        return

    # ----- list_block / list_layout / list_group：items 或 list_items；项内 runs 支持 content.runs -----
    if elem_type in ("list_block", "list_layout", "list_group"):
        items = _get_attr(elem, "items") or _get_attr(elem, "list_items") or []
        first = items[0] if items else None
        if isinstance(first, dict) and _get_attr(first, "type"):
            for child in items:
                _render_element(
                    slide,
                    child,
                    w_inch,
                    h_inch,
                    (left_inch, top_inch),
                    (width_inch, height_inch),
                    icon_images=icon_images,
                )
            return

        tx, ty, tw, th = _inflate_text_region(left_inch, top_inch, width_inch, height_inch, parent_abs)
        shape = slide.shapes.add_textbox(Inches(tx), Inches(ty), Inches(tw), Inches(th))
        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        try:
            tf.margin_left = Pt(DEFAULT_TEXT_PADDING_PT)
            tf.margin_right = Pt(DEFAULT_TEXT_PADDING_PT)
            tf.margin_top = Pt(2)
            tf.margin_bottom = Pt(2)
        except Exception:
            pass

        bullet_type = (_get_attr(elem, "bullet_type") or "disc").lower()
        bullet_char = "•" if bullet_type in {"disc", "dot", "bullet"} else "-"

        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            if hasattr(p, "clear"):
                p.clear()
            p.level = 0
            p.line_spacing = DEFAULT_LINE_SPACING
            p.space_after = Pt(DEFAULT_SPACE_AFTER_PT)
            _force_bullet_paragraph(p, bullet_char=bullet_char)

            if isinstance(item, list):
                runs = item
            else:
                runs = _get_attr(item, "text_runs") or _get_attr(item, "runs")
                if not isinstance(runs, list) and isinstance(_get_attr(item, "content"), dict):
                    runs = _get_attr(_get_attr(item, "content"), "runs") or []
                if not isinstance(runs, list):
                    runs = _get_attr(item, "text") or _get_attr(item, "content") or []
                runs = runs if isinstance(runs, list) else []
            for run_spec in runs:
                run = p.add_run()
                run.text = _run_text(run_spec)
                _apply_run_style(run, run_spec, h_inch)

        if not items:
            tf.paragraphs[0].add_run().text = " "
        return

    # ----- group -----
    if elem_type == "group":
        fill = _get_attr(elem, "background_color") or (
            _get_attr(elem, "background_shape") and _get_attr(_get_attr(elem, "background_shape"), "fill_color")
        )
        border_obj = _get_attr(elem, "border")
        border_color = _get_attr(elem, "border_color") or (
            isinstance(border_obj, dict) and border_obj.get("color")
        ) or (border_obj and _get_attr(border_obj, "color"))
        border_width_raw = _get_attr(elem, "border_width") or (
            isinstance(border_obj, dict) and (border_obj.get("width") or border_obj.get("width_ratio"))
        ) or (border_obj and (_get_attr(border_obj, "width") or _get_attr(border_obj, "width_ratio")))
        border_width_pt = None
        if border_width_raw is not None:
            try:
                v = float(border_width_raw)
                border_width_pt = max(1, int(v * 720)) if 0 < v < 1 else max(1, int(v))
            except (TypeError, ValueError):
                border_width_pt = 2
        if not border_width_pt and border_color:
            border_width_pt = 2
        if fill or border_color:
            _add_rounded_rect(
                slide,
                left_inch,
                top_inch,
                width_inch,
                height_inch,
                fill_color=fill if fill else ("#FFFFFF" if border_color else None),
                border_color=border_color
                or (
                    _get_attr(elem, "background_shape")
                    and _get_attr(_get_attr(elem, "background_shape"), "border")
                    and _get_attr(_get_attr(_get_attr(elem, "background_shape"), "border"), "color")
                ),
                border_width=border_width_pt
                or _get_attr(elem, "border_width")
                or (
                    _get_attr(elem, "background_shape")
                    and _get_attr(_get_attr(elem, "background_shape"), "border")
                    and _get_attr(_get_attr(_get_attr(elem, "background_shape"), "border"), "width")
                ),
                shadow=_get_attr(elem, "shadow")
                if isinstance(_get_attr(elem, "shadow"), bool)
                else (
                    _get_attr(elem, "background_shape")
                    and _get_attr(_get_attr(elem, "background_shape"), "shadow")
                ),
                radius=_get_attr(elem, "border_radius")
                or _get_attr(elem, "radius")
                or (
                    _get_attr(elem, "background_shape")
                    and _get_attr(_get_attr(elem, "background_shape"), "radius")
                ),
                shape_type=_get_attr(elem, "background_shape")
                and _get_attr(_get_attr(elem, "background_shape"), "type"),
            )

        children = _get_attr(elem, "children") or []
        for child in children:
            _render_element(
                slide,
                child,
                w_inch,
                h_inch,
                (left_inch, top_inch),
                (width_inch, height_inch),
                icon_images=icon_images,
            )


def create_slide_v2(
    prs: Presentation,
    doc: Any,
    background: CleanedBackground,
    icon_images: Optional[Dict[str, bytes]] = None,
) -> object:
    """
    按 Schema V2 追加一页：优先铺背景图；若背景净化退化为纯色且文档有 gradient，则用渐变底板提升视觉还原。
    doc 可为 SlideDocumentV2 或 dict：若为 dict 可含 slide_data.elements 或 elements，支持大模型实际输出的 text_box / text_runs / background_color / icon_bbox 等。
    icon_images: placeholder_id -> PNG bytes，用于 icon_text_layout 的图标。None 时会自动注入默认图标（check/lock）。
    """
    if isinstance(doc, dict):
        elements = doc.get("slide_data", {}).get("elements", doc.get("elements", []))
    else:
        elements = getattr(doc, "elements", [])
        if not elements and isinstance(getattr(doc, "slide_data", None), dict):
            elements = doc.slide_data.get("elements", [])

    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    w_inch, h_inch = _slide_dimensions_inch(prs)

    background_drawn = False
    bg_spec = _extract_background_spec(doc)
    if background.method == "solid_fallback":
        background_drawn = _try_draw_gradient_background(slide, w_inch, h_inch, bg_spec)

    if not background_drawn:
        slide.shapes.add_picture(
            io.BytesIO(background.image_bytes),
            Inches(0),
            Inches(0),
            width=Inches(w_inch),
            height=Inches(h_inch),
        )

    final_icon_images = icon_images if icon_images is not None else build_default_icon_images()

    for elem in elements:
        _render_element(
            slide,
            elem,
            w_inch,
            h_inch,
            (0.0, 0.0),
            (w_inch, h_inch),
            icon_images=final_icon_images,
        )
    return slide
