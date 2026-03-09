"""
Pytest 配置与共享 fixtures。

- 临时目录、极小 PDF、固定图片、示例 Slide JSON 等供单元/集成测试使用。
"""
from __future__ import annotations

import io
from pathlib import Path

import fitz
import pytest
from PIL import Image

from app.core.task_store import InMemoryTaskStore, TaskStatus
from app.pipeline.pdf_to_images import PageImage


@pytest.fixture
def tmp_dir(tmp_path: Path) -> Path:
    """临时目录。"""
    return tmp_path


@pytest.fixture
def minimal_pdf_path(tmp_path: Path) -> Path:
    """单页空白 PDF 路径，用于 pdf_to_images / get_pdf_page_count 等测试。"""
    path = tmp_path / "minimal.pdf"
    doc = fitz.open()
    doc.new_page(width=595, height=842)  # A4
    doc.save(str(path))
    doc.close()
    return path


@pytest.fixture
def two_page_pdf_path(tmp_path: Path) -> Path:
    """两页空白 PDF，用于降级测试（如第一页成功、第二页失败）。"""
    path = tmp_path / "two_pages.pdf"
    doc = fitz.open()
    doc.new_page(width=595, height=842)
    doc.new_page(width=595, height=842)
    doc.save(str(path))
    doc.close()
    return path


@pytest.fixture
def sample_page_image() -> PageImage:
    """固定 100x100 红色 PNG 的 PageImage，用于 extract_colors / crop 等。"""
    img = Image.new("RGB", (100, 100), color=(200, 50, 50))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return PageImage(
        image_bytes=buf.getvalue(),
        width=100,
        height=100,
        page_number=1,
    )


SAMPLE_SLIDE_JSON = {
    "id": "slide-1",
    "index": 0,
    "aspectRatio": "16:9",
    "blocks": [
        {
            "id": "text-1",
            "type": "text",
            "content": "Test Title",
            "box": {"x": 0.1, "y": 0.1, "w": 0.8, "h": 0.15},
            "style": {
                "fontSize": 24,
                "fontWeight": "bold",
                "fontColor": "#000000",
                "textAlign": "center",
                "verticalAlign": "middle",
            },
            "zIndex": 1,
        },
        {
            "id": "image-1",
            "type": "image",
            "resourceType": "id",
            "resource": "image-1",
            "box": {"x": 0.2, "y": 0.3, "w": 0.3, "h": 0.3},
            "opacity": 1.0,
            "zIndex": 0,
        },
    ],
    "metadata": {"sourcePage": 1, "mode": "standard", "ocrConfidence": 0.9},
}


@pytest.fixture
def sample_slide_json() -> dict:
    """符合 Slide Schema 的示例 JSON，用于 Gemini mock 或 pptx 合成测试。"""
    return SAMPLE_SLIDE_JSON.copy()


@pytest.fixture
def fresh_task_store() -> InMemoryTaskStore:
    """每个测试独立的 InMemoryTaskStore。"""
    return InMemoryTaskStore()


@pytest.fixture
def completed_task_with_export(tmp_path: Path, fresh_task_store: InMemoryTaskStore) -> tuple[str, Path]:
    """
    在 store 中创建一条已完成且带 exportPath 的任务，并写入一个最小 PPTX 文件。
    返回 (task_id, pptx_path)。
    """
    from pptx import Presentation

    task_id = "test-export-task-id"
    task_dir = tmp_path / task_id
    task_dir.mkdir(parents=True, exist_ok=True)
    pptx_path = task_dir / "output.pptx"
    prs = Presentation()
    prs.slide_width = 12192000  # 13.333 inch
    prs.slide_height = 6858000
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(str(pptx_path))

    fresh_task_store.create_task(
        TaskStatus(
            taskId=task_id,
            filename="test.pptx",
            status="completed",
            progress=100,
            currentPage=0,
            totalPages=1,
            stage="completed",
            failedPages=[],
            exportPath=str(pptx_path),
        )
    )
    return task_id, pptx_path
