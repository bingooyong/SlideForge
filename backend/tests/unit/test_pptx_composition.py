"""PPTX 合成 create_slide、slide_size_for_aspect_ratio 单元测试。"""
from __future__ import annotations

import io

from pptx import Presentation

from app.pipeline.background_cleaning import CleanedBackground
from app.pipeline.color_extraction import StyleInfo
from app.pipeline.layout_ocr_models import Box2D, ImageBlock, Slide, TextBlock
from app.pipeline.icon_cropping import CroppedImage
from app.pipeline.pptx_composition import create_slide, setup_slide_size, slide_size_for_aspect_ratio


def _make_cleaned_background(width: int = 100, height: int = 100) -> CleanedBackground:
    from PIL import Image

    img = Image.new("RGB", (width, height), color=(255, 255, 255))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return CleanedBackground(
        page_number=1,
        image_bytes=buf.getvalue(),
        width=width,
        height=height,
        method="blur",
    )


def _make_slide_schema() -> Slide:
    return Slide(
        id="slide-1",
        index=0,
        aspectRatio="16:9",
        blocks=[
            TextBlock(
                id="text-1",
                content="Hello",
                box=Box2D(x=0.1, y=0.1, w=0.5, h=0.1),
            ),
            ImageBlock(
                id="img-1",
                resourceType="id",
                resource="img-1",
                box=Box2D(x=0.2, y=0.3, w=0.3, h=0.3),
            ),
        ],
        metadata=None,
    )


def _make_cropped_image() -> CroppedImage:
    from PIL import Image

    img = Image.new("RGB", (30, 30), color=(100, 100, 100))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return CroppedImage(
        block_id="img-1",
        page_number=1,
        image_bytes=buf.getvalue(),
        width=30,
        height=30,
    )


def test_slide_size_for_aspect_ratio() -> None:
    """slide_size_for_aspect_ratio 应返回正确英寸尺寸。"""
    w, h = slide_size_for_aspect_ratio("16:9")
    assert w == 13.333
    assert h == 7.5

    w, h = slide_size_for_aspect_ratio("4:3")
    assert w == 10.0
    assert h == 7.5

    w, h = slide_size_for_aspect_ratio("unknown")
    assert w == 13.333
    assert h == 7.5


def test_create_slide_adds_slide_to_presentation() -> None:
    """create_slide 应在 prs 中追加一页。"""
    prs = Presentation()
    setup_slide_size(prs, "16:9")
    initial_count = len(prs.slides)

    create_slide(
        prs=prs,
        slide_schema=_make_slide_schema(),
        background=_make_cleaned_background(),
        cropped_images=[_make_cropped_image()],
        style_info=StyleInfo(
            primaryColor="#000000",
            backgroundColor="#FFFFFF",
            accentColors=["#000000"],
        ),
    )

    assert len(prs.slides) == initial_count + 1
